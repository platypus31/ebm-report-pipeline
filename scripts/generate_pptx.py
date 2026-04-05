#!/usr/bin/env python3
"""
EBM Report PowerPoint Generator (Fallback)
當 Canva MCP 不可用時，使用 python-pptx 產生可編輯的 .pptx 檔案。

用法：
    python3 scripts/generate_pptx.py slides.json output.pptx

slides.json 格式：
{
    "title": "EBM 報告標題",
    "author": "報告者",
    "department": "科別",
    "date": "2026-04-04",
    "style": "formal",  // formal | clean | teaching | competition
    "slides": [
        {
            "type": "title",
            "title": "標題",
            "subtitle": "副標題"
        },
        {
            "type": "toc",
            "title": "目錄",
            "items": ["ASK 問題", "ACQUIRE 檢索", "APPRAISE 評讀", "APPLY 應用", "AUDIT 評估"]
        },
        {
            "type": "section",
            "title": "ASK",
            "subtitle": "問題"
        },
        {
            "type": "content",
            "title": "投影片標題",
            "bullets": ["重點一", "重點二", "重點三"],
            "section": "ASK"
        },
        {
            "type": "two_column",
            "title": "標題",
            "left_title": "左欄標題",
            "left_bullets": ["..."],
            "right_title": "右欄標題",
            "right_bullets": ["..."],
            "section": "APPRAISE"
        },
        {
            "type": "table",
            "title": "標題",
            "headers": ["欄1", "欄2", "欄3"],
            "rows": [["值1", "值2", "值3"]],
            "section": "ACQUIRE"
        },
        {
            "type": "image",
            "title": "截圖標題",
            "image_path": "assets/screenshots/pubmed-search-1.png",
            "caption": "PubMed 搜尋結果",
            "section": "ACQUIRE"
        },
        {
            "type": "image_content",
            "title": "標題",
            "image_path": "assets/screenshots/article-results-1.png",
            "bullets": ["重點一", "重點二"],
            "section": "APPRAISE"
        },
        {
            "type": "appraisal",
            "title": "Section A — Validity",
            "questions": [
                {"number": "Q1", "text": "問題文字", "answer": "Yes", "evidence": "佐證"},
                {"number": "Q2", "text": "問題文字", "answer": "No", "evidence": "佐證"},
                {"number": "Q3", "text": "問題文字", "answer": "Can't tell", "evidence": "佐證"}
            ],
            "section": "APPRAISE"
        }
    ]
}
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:
    print("錯誤：缺少 python-pptx 套件。請執行 pip install python-pptx 安裝。", file=sys.stderr)
    sys.exit(1)

# ── Constants ──

VALID_STYLES = ["formal", "clean", "teaching", "competition"]
VALID_SLIDE_TYPES = [
    "title", "section", "content", "two_column", "table",
    "image", "image_content", "toc", "appraisal",
]

# ── Style definitions ──

STYLES = {
    "formal": {
        "primary": RGBColor(0x1B, 0x3A, 0x5C),      # 深藍
        "secondary": RGBColor(0x2E, 0x86, 0xC1),     # 中藍
        "accent": RGBColor(0x17, 0xA5, 0x89),         # 青綠
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "light": RGBColor(0xEB, 0xF5, 0xFB),
    },
    "clean": {
        "primary": RGBColor(0x34, 0x49, 0x5E),
        "secondary": RGBColor(0x00, 0xB8, 0x94),      # 綠色
        "accent": RGBColor(0xF3, 0x9C, 0x12),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "light": RGBColor(0xF0, 0xF9, 0xF4),
    },
    "teaching": {
        "primary": RGBColor(0xE6, 0x7E, 0x22),       # 橘色
        "secondary": RGBColor(0x29, 0x80, 0xB9),
        "accent": RGBColor(0x27, 0xAE, 0x60),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "light": RGBColor(0xFD, 0xF2, 0xE9),
    },
    "competition": {
        "primary": RGBColor(0x1A, 0x23, 0x5B),       # 深藍紫
        "secondary": RGBColor(0x3D, 0x5A, 0xFE),
        "accent": RGBColor(0x00, 0xE6, 0x76),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1A, 0x1A, 0x2E),
        "light": RGBColor(0xE8, 0xEA, 0xF6),
    },
}

# Appraisal answer colors
ANSWER_COLORS = {
    "Yes": RGBColor(0x27, 0xAE, 0x60),       # 綠色
    "No": RGBColor(0xE7, 0x4C, 0x3C),         # 紅色
    "Can't tell": RGBColor(0x95, 0xA5, 0xA6),  # 灰色
    "N/A": RGBColor(0x7F, 0x8C, 0x8D),         # 深灰
}

SECTION_LABELS = {
    "ASK": "問題 ASK",
    "ACQUIRE": "檢索 ACQUIRE",
    "APPRAISE": "評讀 APPRAISE",
    "APPLY": "應用 APPLY",
    "AUDIT": "評估 AUDIT",
}


def add_section_indicator(slide, current_section, style_colors):
    """Add 5A section indicator bar on the left side."""
    sections = ["ASK", "ACQUIRE", "APPRAISE", "APPLY", "AUDIT"]
    bar_top = Inches(1.5)
    bar_height = Inches(0.6)
    bar_width = Inches(1.2)
    bar_left = Inches(0.15)

    for i, sec in enumerate(sections):
        top = bar_top + Emu(int(bar_height * i * 1.1))
        is_current = (sec == current_section)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_left, top, bar_width, bar_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = style_colors["primary"] if is_current else style_colors["light"]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = SECTION_LABELS[sec]
        p.font.size = Pt(8)
        p.font.bold = is_current
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if is_current else style_colors["text"]
        p.alignment = PP_ALIGN.CENTER


def create_title_slide(prs, data, style_colors):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style_colors["primary"]
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "EBM Report")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = style_colors["primary"]

    # Author / Department / Date
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    lines = []
    if data.get("author"):
        lines.append(data["author"])
    if data.get("department"):
        lines.append(data["department"])
    if data.get("date"):
        lines.append(data["date"])
    # Advisor support
    if data.get("advisor"):
        lines.append(f"指導老師：{data['advisor']}")
    p = tf2.paragraphs[0]
    p.text = " | ".join(lines)
    p.font.size = Pt(18)
    p.font.color.rgb = style_colors["secondary"]


def create_toc_slide(prs, slide_data, style_colors):
    """Create Table of Contents / 5A overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    _add_title_bar(slide, slide_data.get("title", "目錄 / 5A 流程"), Inches(0.8), style_colors)

    items = slide_data.get("items", [
        "ASK — 提出臨床問題",
        "ACQUIRE — 搜尋最佳證據",
        "APPRAISE — 嚴格評讀",
        "APPLY — 應用到臨床",
        "AUDIT — 自我評估",
    ])

    sections_5a = ["ASK", "ACQUIRE", "APPRAISE", "APPLY", "AUDIT"]
    box_width = Inches(1.5)
    box_height = Inches(1.2)
    arrow_width = Inches(0.4)
    total_width = len(sections_5a) * box_width + (len(sections_5a) - 1) * arrow_width
    start_left = Inches((10 - total_width / Inches(1)) / 2)

    # Draw 5A flow boxes
    for i, sec in enumerate(sections_5a):
        left = start_left + Emu(int((box_width + arrow_width) * i))
        top = Inches(2.5)

        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = style_colors["primary"]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sec
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # Arrow between boxes (except last)
        if i < len(sections_5a) - 1:
            arrow_left = left + box_width
            arrow_top = top + Emu(int(box_height * 0.35))
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left, arrow_top, arrow_width, Emu(int(box_height * 0.3))
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = style_colors["accent"]
            arrow.line.fill.background()

    # Items below
    if items:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = style_colors["text"]
            p.space_after = Pt(6)


def create_section_slide(prs, slide_data, style_colors):
    """Create 5A section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style_colors["primary"]
    shape.line.fill.background()

    # Section title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if slide_data.get("subtitle"):
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = slide_data["subtitle"]
        p2.font.size = Pt(24)
        p2.font.color.rgb = style_colors["accent"]
        p2.alignment = PP_ALIGN.CENTER


def _add_title_bar(slide, title, content_left, style_colors):
    """Add a title bar rectangle and title text to a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style_colors["primary"]
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(content_left, Inches(0.12), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def create_content_slide(prs, slide_data, style_colors):
    """Create standard content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    section = slide_data.get("section")
    content_left = Inches(1.6) if section else Inches(0.8)

    # Section indicator
    if section:
        add_section_indicator(slide, section, style_colors)

    _add_title_bar(slide, slide_data.get("title", ""), content_left, style_colors)

    # Bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        txBox2 = slide.shapes.add_textbox(content_left, Inches(1.2), Inches(8), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"  {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = style_colors["text"]
            p.space_after = Pt(8)
            p.level = 0


def create_two_column_slide(prs, slide_data, style_colors):
    """Create two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    section = slide_data.get("section")
    content_left = Inches(1.6) if section else Inches(0.5)

    if section:
        add_section_indicator(slide, section, style_colors)

    _add_title_bar(slide, slide_data.get("title", ""), content_left, style_colors)

    col_width = Inches(3.8) if section else Inches(4.2)

    # Left column
    left_box = slide.shapes.add_textbox(content_left, Inches(1.2), col_width, Inches(5.5))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = slide_data.get("left_title", "")
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = style_colors["secondary"]
    for bullet in slide_data.get("left_bullets", []):
        p = tf_l.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(14)
        p.font.color.rgb = style_colors["text"]
        p.space_after = Pt(6)

    # Right column
    right_left = content_left + col_width + Inches(0.3)
    right_box = slide.shapes.add_textbox(right_left, Inches(1.2), col_width, Inches(5.5))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = slide_data.get("right_title", "")
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = style_colors["secondary"]
    for bullet in slide_data.get("right_bullets", []):
        p = tf_r.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(14)
        p.font.color.rgb = style_colors["text"]
        p.space_after = Pt(6)


def create_table_slide(prs, slide_data, style_colors):
    """Create slide with table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    section = slide_data.get("section")
    content_left = Inches(1.6) if section else Inches(0.5)

    if section:
        add_section_indicator(slide, section, style_colors)

    _add_title_bar(slide, slide_data.get("title", ""), content_left, style_colors)

    headers = slide_data.get("headers", [])
    rows_data = slide_data.get("rows", [])

    if not headers:
        return

    n_rows = len(rows_data) + 1
    n_cols = len(headers)
    table_width = Inches(7.5) if section else Inches(9)
    row_height = Inches(0.5)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        content_left, Inches(1.2),
        table_width, Emu(int(row_height * n_rows))
    )
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = style_colors["primary"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = style_colors["light"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = style_colors["text"]


def create_image_slide(prs, slide_data, style_colors, project_dir=None):
    """Create full-page image slide with title and optional caption."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    section = slide_data.get("section")
    content_left = Inches(1.6) if section else Inches(0.5)

    if section:
        add_section_indicator(slide, section, style_colors)

    _add_title_bar(slide, slide_data.get("title", ""), content_left, style_colors)

    # Try to insert image
    image_path = slide_data.get("image_path", "")
    image_inserted = False

    if image_path and project_dir:
        full_path = Path(project_dir) / image_path
        if full_path.exists():
            img_left = content_left + Inches(0.3)
            img_top = Inches(1.3)
            img_width = Inches(7.0) if section else Inches(8.5)
            img_height = Inches(5.0)
            try:
                slide.shapes.add_picture(
                    str(full_path), img_left, img_top, img_width, img_height
                )
                image_inserted = True
            except Exception:
                pass

    if not image_inserted:
        # Placeholder box
        ph_left = content_left + Inches(0.5)
        ph_top = Inches(1.5)
        ph_width = Inches(7.0) if section else Inches(8.5)
        ph_height = Inches(4.5)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, ph_left, ph_top, ph_width, ph_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = style_colors["light"]
        shape.line.color.rgb = style_colors["secondary"]

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        display_path = image_path or "（未指定圖片路徑）"
        p.text = f"📷 {display_path}"
        p.font.size = Pt(14)
        p.font.color.rgb = style_colors["secondary"]
        p.alignment = PP_ALIGN.CENTER

    # Caption
    caption = slide_data.get("caption", "")
    if caption:
        cap_box = slide.shapes.add_textbox(
            content_left, Inches(6.5), Inches(8), Inches(0.6)
        )
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = style_colors["secondary"]
        p.alignment = PP_ALIGN.CENTER


def create_image_content_slide(prs, slide_data, style_colors, project_dir=None):
    """Create slide with image on left and bullet points on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    section = slide_data.get("section")
    content_left = Inches(1.6) if section else Inches(0.5)

    if section:
        add_section_indicator(slide, section, style_colors)

    _add_title_bar(slide, slide_data.get("title", ""), content_left, style_colors)

    # Image area (left half)
    img_width = Inches(4.0) if section else Inches(4.5)
    image_path = slide_data.get("image_path", "")
    image_inserted = False

    if image_path and project_dir:
        full_path = Path(project_dir) / image_path
        if full_path.exists():
            try:
                slide.shapes.add_picture(
                    str(full_path), content_left, Inches(1.3),
                    img_width, Inches(5.0)
                )
                image_inserted = True
            except Exception:
                pass

    if not image_inserted:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            content_left, Inches(1.3), img_width, Inches(5.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = style_colors["light"]
        shape.line.color.rgb = style_colors["secondary"]
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        display_path = image_path or "（未指定圖片路徑）"
        p.text = f"📷 {display_path}"
        p.font.size = Pt(12)
        p.font.color.rgb = style_colors["secondary"]
        p.alignment = PP_ALIGN.CENTER

    # Bullets area (right half)
    bullets_left = content_left + img_width + Inches(0.3)
    bullets_width = Inches(3.5) if section else Inches(4.2)
    bullets = slide_data.get("bullets", [])

    if bullets:
        txBox = slide.shapes.add_textbox(
            bullets_left, Inches(1.3), bullets_width, Inches(5.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {bullet}"
            p.font.size = Pt(14)
            p.font.color.rgb = style_colors["text"]
            p.space_after = Pt(6)


def _get_answer_color(answer):
    """Return color for appraisal answer."""
    answer_lower = answer.strip().lower()
    if answer_lower == "yes":
        return ANSWER_COLORS["Yes"]
    elif answer_lower == "no":
        return ANSWER_COLORS["No"]
    elif answer_lower in ("can't tell", "cant tell", "unclear"):
        return ANSWER_COLORS["Can't tell"]
    else:
        return ANSWER_COLORS["N/A"]


def create_appraisal_slide(prs, slide_data, style_colors):
    """Create appraisal slide with color-coded answers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    section = slide_data.get("section", "APPRAISE")
    content_left = Inches(1.6) if section else Inches(0.5)

    if section:
        add_section_indicator(slide, section, style_colors)

    _add_title_bar(slide, slide_data.get("title", ""), content_left, style_colors)

    questions = slide_data.get("questions", [])
    if not questions:
        return

    # Build table: number | question | answer | evidence
    show_evidence = any(q.get("evidence") for q in questions)
    if show_evidence:
        headers = ["題號", "問題", "判定", "佐證"]
        n_cols = 4
        col_widths = [Inches(0.6), Inches(3.0), Inches(0.8), Inches(3.1)]
    else:
        headers = ["題號", "問題", "判定"]
        n_cols = 3
        col_widths = [Inches(0.7), Inches(4.5), Inches(1.0)]

    n_rows = len(questions) + 1
    table_width = sum(w for w in col_widths)
    row_height = Inches(0.55)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        content_left, Inches(1.2),
        table_width, Emu(int(row_height * n_rows))
    )
    table = table_shape.table

    # Set column widths
    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = style_colors["primary"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, q in enumerate(questions):
        answer = q.get("answer", "")
        answer_color = _get_answer_color(answer)

        # Number
        cell_num = table.cell(i + 1, 0)
        cell_num.text = q.get("number", f"Q{i + 1}")
        for p in cell_num.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = style_colors["text"]
            p.alignment = PP_ALIGN.CENTER

        # Question text
        cell_q = table.cell(i + 1, 1)
        cell_q.text = q.get("text", "")
        for p in cell_q.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = style_colors["text"]

        # Answer (color-coded)
        cell_a = table.cell(i + 1, 2)
        cell_a.text = answer
        cell_a.fill.solid()
        cell_a.fill.fore_color.rgb = answer_color
        for p in cell_a.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

        # Evidence (if shown)
        if show_evidence:
            cell_e = table.cell(i + 1, 3)
            cell_e.text = q.get("evidence", "")
            for p in cell_e.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = style_colors["text"]

        # Alternating row color (except answer column)
        if i % 2 == 0:
            for j in [0, 1] + ([3] if show_evidence else []):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = style_colors["light"]


def validate_data(data):
    """Validate input JSON data structure. Raises ValueError on invalid input."""
    if not isinstance(data, dict):
        raise ValueError("輸入資料必須是 JSON 物件（字典）。")

    # Required fields
    missing = []
    for field in ["title", "style", "slides"]:
        if field not in data:
            missing.append(field)
    if missing:
        raise ValueError(f"缺少必要欄位：{', '.join(missing)}。JSON 必須包含 title、style、slides。")

    # Validate style
    style = data["style"]
    if style not in VALID_STYLES:
        raise ValueError(
            f"無效的簡報風格：'{style}'。"
            f"可用的風格為：{', '.join(VALID_STYLES)}。"
        )

    # Validate slides
    slides = data["slides"]
    if not isinstance(slides, list):
        raise ValueError("'slides' 欄位必須是陣列。")
    if len(slides) == 0:
        raise ValueError("'slides' 陣列不可為空，至少需要一張投影片。")

    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            raise ValueError(f"第 {i + 1} 張投影片必須是 JSON 物件。")
        slide_type = slide.get("type")
        if slide_type is None:
            raise ValueError(f"第 {i + 1} 張投影片缺少 'type' 欄位。")
        if slide_type not in VALID_SLIDE_TYPES:
            raise ValueError(
                f"第 {i + 1} 張投影片的類型 '{slide_type}' 無效。"
                f"可用的類型為：{', '.join(VALID_SLIDE_TYPES)}。"
            )


def generate_pptx(data, output_path, project_dir=None):
    """Generate the PowerPoint file from structured data."""
    validate_data(data)

    style_name = data.get("style", "formal")
    style_colors = STYLES.get(style_name, STYLES["formal"])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create title slide from metadata (unless slides array starts with a title slide)
    slides_list = data.get("slides", [])
    first_is_title = slides_list and slides_list[0].get("type") == "title"
    if not first_is_title:
        create_title_slide(prs, data, style_colors)

    # Create slides
    for slide_data in slides_list:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            create_title_slide(prs, {**data, **slide_data}, style_colors)
        elif slide_type == "toc":
            create_toc_slide(prs, slide_data, style_colors)
        elif slide_type == "section":
            create_section_slide(prs, slide_data, style_colors)
        elif slide_type == "content":
            create_content_slide(prs, slide_data, style_colors)
        elif slide_type == "two_column":
            create_two_column_slide(prs, slide_data, style_colors)
        elif slide_type == "table":
            create_table_slide(prs, slide_data, style_colors)
        elif slide_type == "image":
            create_image_slide(prs, slide_data, style_colors, project_dir)
        elif slide_type == "image_content":
            create_image_content_slide(prs, slide_data, style_colors, project_dir)
        elif slide_type == "appraisal":
            create_appraisal_slide(prs, slide_data, style_colors)
        else:
            create_content_slide(prs, slide_data, style_colors)

    try:
        prs.save(output_path)
    except OSError as e:
        raise OSError(f"無法儲存 PowerPoint 至 '{output_path}'：{e}") from e
    print(f"PowerPoint 已儲存：{output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="EBM Report PowerPoint 產生器 — 將 JSON 簡報資料轉換為 .pptx 檔案",
        epilog=(
            "範例：\n"
            "  python3 generate_pptx.py slides.json output.pptx\n"
            "  python3 generate_pptx.py slides.json ~/Desktop/ebm-report.pptx\n"
            "  python3 generate_pptx.py slides.json output.pptx --project projects/my-topic\n\n"
            "JSON 格式說明：\n"
            "  必要欄位：title (字串), style (formal|clean|teaching|competition), slides (陣列)\n"
            "  每張投影片的 type 可為：title, toc, section, content, two_column, table,\n"
            "                         image, image_content, appraisal\n"
            "  詳見 skills/ebm-slides.md 中的 JSON 範例。"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("input", help="輸入的 JSON 檔案路徑（簡報資料）")
    parser.add_argument("output", help="輸出的 .pptx 檔案路徑")
    parser.add_argument("--project", help="專案目錄路徑（用於解析截圖圖片）")

    args = parser.parse_args()
    input_path = args.input
    output_path = args.output

    # Validate input file exists
    if not Path(input_path).is_file():
        print(f"錯誤：找不到輸入檔案 '{input_path}'。請確認路徑是否正確。", file=sys.stderr)
        sys.exit(1)

    # Load and parse JSON
    try:
        with open(input_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"錯誤：JSON 格式不正確。解析失敗於第 {e.lineno} 行，第 {e.colno} 欄。", file=sys.stderr)
        print(f"詳細資訊：{e.msg}", file=sys.stderr)
        sys.exit(1)
    except UnicodeDecodeError:
        print("錯誤：檔案編碼不正確。請確認檔案為 UTF-8 編碼。", file=sys.stderr)
        sys.exit(1)

    # Validate and generate
    try:
        generate_pptx(data, output_path, project_dir=args.project)
    except ValueError as e:
        print(f"錯誤：輸入資料驗證失敗。{e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"錯誤：產生簡報時發生未預期的錯誤。{e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
